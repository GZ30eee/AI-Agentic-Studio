import os
import hashlib
import time
import json
import logging
import sqlite3
from datetime import datetime
from langgraph.graph import StateGraph, END
from typing import TypedDict, Type, List, Dict, Any, Optional
import threading
from queue import Queue

import smtplib
import logging
# Duplicate logger removed

from langgraph.graph import StateGraph
import streamlit as st
from streamlit_option_menu import option_menu
import streamlit_authenticator as stauth
import yaml
from yaml.loader import SafeLoader

from tenacity import retry, stop_after_attempt, wait_fixed
from dotenv import load_dotenv
from typing import cast, TypedDict, Type, List, Dict, Any, Optional

# Core Frameworks
from crewai import Agent, Task, Crew, LLM, Process
from crewai.tools import BaseTool
from langchain_community.tools import DuckDuckGoSearchRun
from crewai.tools import BaseTool
from pydantic import BaseModel, Field

os.environ["CREWAI_TELEMETRY_OPT_OUT"] = "true"

# PDF & Email
from reportlab.lib.pagesizes import LETTER
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.application import MIMEApplication

# New imports for enhancements
import requests
from bs4 import BeautifulSoup
import textstat
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import markdown
import chromadb
from chromadb.utils import embedding_functions
import sqlalchemy as sa
from sqlalchemy.orm import sessionmaker, declarative_base
import openai
import anthropic
import google.generativeai as genai

# Additional imports for new features
import io
import csv
import PyPDF2
from gtts import gTTS
from streamlit.components.v1 import html as st_html
from xml.sax.saxutils import escape as xml_escape

from pydantic import BaseModel, Field

class SearchInput(BaseModel):
    query: str = Field(..., description="Search query")

# For monitoring (stub)
from prometheus_client import Counter, Histogram, start_http_server

load_dotenv()

# --- LOGGING ---
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# --- MONITORING SETUP (optional) ---
if os.getenv("ENABLE_MONITORING", "false").lower() == "true":
    start_http_server(8000)
    REQUEST_COUNT = Counter('app_requests_total', 'Total requests')
    REQUEST_DURATION = Histogram('app_request_duration_seconds', 'Request duration')
else:
    class DummyCounter:
        def inc(self): pass
    class DummyHistogram:
        def time(self): return self
        def __enter__(self): pass
        def __exit__(self, *args): pass
    REQUEST_COUNT = DummyCounter()
    REQUEST_DURATION = DummyHistogram()

# --- DATABASE SETUP (SQLite) ---
engine = sa.create_engine('sqlite:///reports.db')
SessionLocal = sessionmaker(bind=engine)
Base = declarative_base()

class ReportDB(Base):
    __tablename__ = 'reports'
    id = sa.Column(sa.Integer, primary_key=True)
    topic = sa.Column(sa.String)
    email = sa.Column(sa.String)
    report_text = sa.Column(sa.Text)
    pdf_path = sa.Column(sa.String)
    created_at = sa.Column(sa.DateTime, default=datetime.utcnow)
    user_id = sa.Column(sa.String)

Base.metadata.create_all(engine)

# --- TOOL DEFINITIONS ---
class SearchInput(BaseModel):
    query: str = Field(..., description="Search query")

class DuckDuckGoTool(BaseTool):
    name: str = "DuckDuckGo Search"
    description: str = "Search the web using DuckDuckGo"
    args_schema: Type[BaseModel] = SearchInput

    def _run(self, query: str) -> str:
        search = DuckDuckGoSearchRun()
        return search.run(query)

class WebScraperInput(BaseModel):
    url: str = Field(..., description="URL to scrape")

class WebScraperTool(BaseTool):
    name: str = "Web Scraper"
    description: str = "Fetch and extract text content from a given URL"
    args_schema: Type[BaseModel] = WebScraperInput

    def _run(self, url: str) -> str:
        try:
            response = requests.get(url, timeout=10)
            soup = BeautifulSoup(response.text, 'html.parser')
            for script in soup(["script", "style"]):
                script.decompose()
            text = soup.get_text()
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            return text[:5000]
        except Exception as e:
            return f"Scraping error: {e}"

class NewsAPIInput(BaseModel):
    query: str = Field(..., description="Search query for news")

class NewsAPITool(BaseTool):
    name: str = "News API"
    description: str = "Fetch recent news articles on a topic"
    args_schema: Type[BaseModel] = NewsAPIInput

    def _run(self, query: str) -> str:
        api_key = os.getenv("NEWS_API_KEY")
        if not api_key:
            return "News API key not configured."
        url = f"https://newsapi.org/v2/everything?q={query}&apiKey={api_key}"
        try:
            response = requests.get(url).json()
            articles = response.get('articles', [])[:5]
            summaries = []
            for art in articles:
                summaries.append(f"Title: {art['title']}\nDescription: {art['description']}\nURL: {art['url']}")
            return "\n\n".join(summaries)
        except Exception as e:
            return f"News API error: {e}"

class RAGTool(BaseTool):
    name: str = "Company Documents RAG"
    description: str = "Retrieve information from internal documents"
    args_schema: Type[BaseModel] = SearchInput
    
    class Config:
        arbitrary_types_allowed = True
    
    client: Any = Field(description="ChromaDB client", default=None)
    embedding_func: Any = Field(description="Embedding function", default=None)
    collection: Any = Field(description="ChromaDB collection", default=None)

    def __init__(self, collection_name="company_docs", **kwargs):
        super().__init__(**kwargs)
        try:
            client = chromadb.PersistentClient(path="./chroma_db")
        except Exception as e:
            logger.warning(f"Persistent ChromaDB failed, using ephemeral client: {e}")
            client = chromadb.EphemeralClient()
        
        embedding_func = embedding_functions.SentenceTransformerEmbeddingFunction(
            model_name="all-MiniLM-L6-v2"
        )
        
        collection = client.get_or_create_collection(
            name=collection_name, embedding_function=embedding_func
        )
        
        object.__setattr__(self, "client", client)
        object.__setattr__(self, "embedding_func", embedding_func)
        object.__setattr__(self, "collection", collection)

    def add_documents(self, texts, metadatas=None):
        ids = [hashlib.md5(text.encode()).hexdigest() for text in texts]
        self.collection.add(documents=texts, metadatas=metadatas, ids=ids)

    def _run(self, query: str) -> str:
        results = self.collection.query(query_texts=[query], n_results=5)
        docs = results['documents'][0]
        return "\n\n".join(docs) if docs else "No relevant documents found."

# Instantiate tools
search_tool = DuckDuckGoTool()
scraper_tool = WebScraperTool()
news_tool = NewsAPITool()
# rag_tool will be instantiated inside the app to allow session isolation
rag_tool = None 

# --- FREE MODELS DICTIONARY ---
FREE_MODELS = {
    "OpenAI": ["gpt-3.5-turbo", "gpt-4o-mini"],
    "Anthropic": ["claude-3-haiku-20240307"],
    "Gemini": ["gemini-1.5-flash"],
    "Ollama": ["phi3:mini"]
}

# --- LLM FACTORY (unchanged) ---
class LLMFactory:
    @staticmethod
    def get_llm(provider: str, model: str, temperature: float = 0.7, max_tokens: int = 4096):
        if provider == "Ollama":
            return LLM(
                model=f"ollama/{model}",
                base_url=os.getenv("OLLAMA_URL", "http://localhost:11434"),
                temperature=temperature,
                max_tokens=max_tokens
            )
        elif provider == "OpenAI":
            return LLM(
                model=f"openai/{model}",
                api_key=os.getenv("OPENAI_API_KEY"),
                temperature=temperature,
                max_tokens=max_tokens
            )
        elif provider == "Anthropic":
            return LLM(
                model=f"anthropic/{model}",
                api_key=os.getenv("ANTHROPIC_API_KEY"),
                temperature=temperature,
                max_tokens=max_tokens
            )
        elif provider == "Gemini":
            return LLM(
                model=f"gemini/{model}",
                api_key=os.getenv("GOOGLE_API_KEY"),
                temperature=temperature,
                max_tokens=max_tokens
            )
        else:
            raise ValueError(f"Unsupported provider: {provider}")

# --- STATE DEFINITION ---
class AgentState(TypedDict):
    topic: str
    email: str
    plan: str
    raw_research: str
    final_report: str
    quality_score: float
    status: str
    num_agents: int
    agent_roles: List[str]
    writing_style: str
    llm_provider: str
    llm_model: str
    citations: List[Dict[str, str]]
    research_notes: List[str]
    fact_check_report: str
    readability_scores: Dict[str, float]
    target_language: Optional[str]
    cache_key: str
    refinement_attempts: int

# --- CACHE ---
CACHE = {}
# Removed unused get_cached_or_search

# --- NODES (unchanged) ---
@retry(stop=stop_after_attempt(3), wait=wait_fixed(2))
def planning_node(state: AgentState):
    logger.info(f"Planning: {state['topic']}")
    llm = LLMFactory.get_llm(state.get('llm_provider', 'Ollama'), state.get('llm_model', 'llama3.1'))
    director = Agent(
        role='Research Director',
        goal=f'Strategic plan for {state["topic"]}',
        backstory="Senior strategist specialized in deep-dive whitepapers.",
        llm=llm,
        max_rpm=50
    )
    plan_prompt = f"Plan an extensive {state.get('num_agents',5)}-section research structure for {state['topic']} with a {state.get('writing_style','formal')} tone."
    task = Task(description=plan_prompt, agent=director, expected_output="A detailed roadmap.")
    result_obj = Crew(agents=[director], tasks=[task], timeout=120).kickoff()
    return {"plan": str(result_obj.raw), "status": "Structure Planned"}

@retry(stop=stop_after_attempt(3), wait=wait_fixed(2))
def research_node(state: AgentState):
    agents = []
    tasks = []
    llm = LLMFactory.get_llm(state['llm_provider'], state['llm_model'])
    plan_str = str(state.get('plan', ''))
    num_val = state.get('num_agents', 1)
    num_agents = int(num_val) if num_val is not None else 1
    
    for i in range(num_agents):
        researcher = Agent(
            role=f'Researcher {i+1}',
            goal=f'Execute part {i+1} of the research plan. Plan snippet: {plan_str[:500]}',
            backstory=f"Expert analyst focusing on sub-aspects of {state['topic']}.",
            tools=[search_tool, scraper_tool, news_tool, rag_tool] if (rag_tool is not None) else [search_tool, scraper_tool, news_tool],
            llm=llm,
            max_rpm=50,
            verbose=True
        )
        agents.append(researcher)
        task = Task(
            description=f"Gather detailed information on {state['topic']}. Focus on sub-aspect {i+1}.",
            agent=researcher,
            expected_output="Detailed research notes with citations."
        )
        tasks.append(task)

    crew = Crew(
        agents=agents, 
        tasks=tasks, 
        process=Process.sequential,
        verbose=True
    )
    
    result_obj = crew.kickoff()
    combined_notes = "\n\n".join([out.raw for out in result_obj.tasks_output])
    
    return {
        "raw_research": combined_notes, 
        "research_notes": [out.raw for out in result_obj.tasks_output], 
        "citations": [], 
        "status": "Parallel Research Complete"
    }

@retry(stop=stop_after_attempt(3), wait=wait_fixed(2))
def writing_node(state: AgentState):
    llm = LLMFactory.get_llm(state['llm_provider'], state['llm_model'])
    style_guide = {
        'formal': "Use a professional, formal tone with precise language.",
        'persuasive': "Use persuasive language to convince executives.",
        'technical': "Include technical details, data, and code snippets where relevant."
    }.get(state['writing_style'], "Use a professional tone.")

    writer = Agent(
        role='Senior Technical Consultant',
        goal=f'Create a high-impact, long-form executive whitepaper in {state["writing_style"]} style.',
        backstory="Lead consultant.",
        llm=llm,
        max_rpm=50
    )
    task = Task(
        description=f"""
Based on this research: {state.get('raw_research', '')},
write a comprehensive technical whitepaper with the following sections:
1. EXECUTIVE SUMMARY
2. MARKET DYNAMICS & RECENT TRENDS
3. STRATEGIC RECOMMENDATIONS FOR ENTERPRISES
4. TECHNICAL CHALLENGES & SOLUTIONS
5. FUTURE OUTLOOK & CONCLUSION

Writing style: {style_guide}
Be verbose and aim for at least 1500 words. Include citations where appropriate.

**IMPORTANT**: Include at least one Mermaid.js diagram illustrating the core concepts, wrapped in ```mermaid tags. The diagram should be placed in a relevant section.
""",
        agent=writer,
        expected_output="A multi-page detailed report with Mermaid diagrams."
    )
    result = Crew(agents=[writer], tasks=[task], timeout=300).kickoff()
    attempts_val = state.get('refinement_attempts', 0)
    attempts = (int(attempts_val) if attempts_val is not None else 0) + 1
    return {"final_report": result.raw, "status": "Report Generated", "refinement_attempts": attempts}

def quality_node(state: AgentState):
    report = state.get("final_report", "")
    length_score = len(report) / 800
    flesch = textstat.flesch_reading_ease(report)
    gunning_fog = textstat.gunning_fog(report)
    sentences = report.split('. ')
    complex_sentences = [s for s in sentences if len(s.split()) > 30]
    readability = {
        "flesch_reading_ease": flesch,
        "gunning_fog": gunning_fog,
        "complex_sentences_count": len(complex_sentences)
    }
    quality_score = length_score * 0.5 + (flesch/100) * 0.5
    return {"quality_score": quality_score, "readability_scores": readability, "status": "Quality Checked"}

def fact_check_node(state: AgentState):
    report = state.get("final_report", "")
    llm = LLMFactory.get_llm(state['llm_provider'], state['llm_model'])
    fact_checker = Agent(
        role='Fact Checker',
        goal='Verify factual accuracy of the report',
        backstory="You are a meticulous researcher who checks every claim against known facts.",
        llm=llm
    )
    task = Task(
        description=f"Review the following report for factual inaccuracies. List any questionable claims and suggest corrections:\n\n{report}",
        agent=fact_checker,
        expected_output="A list of potential inaccuracies and corrections."
    )
    result = Crew(agents=[fact_checker], tasks=[task]).kickoff()
    return {"fact_check_report": result.raw, "status": "Fact Checked"}

def citation_node(state: AgentState):
    citations = cast(list, state.get("citations", []))
    if not isinstance(citations, list):
        citations = []
    
    if not citations:
        citations = [] # Removed dummy data
    
    report = str(state.get("final_report", ""))
    ref_section = "\n\n## References\n"
    for i, cit in enumerate(citations, 1):
        cit_dict = cast(dict, cit)
        source = str(cit_dict.get('source', 'Unknown'))
        snippet = str(cit_dict.get('snippet', ''))
        ref_section += f"{i}. {source} - {snippet[:100]}...\n"
    updated_report = report + ref_section
    return {"final_report": updated_report, "citations": citations, "status": "Citations Added"}

def translation_node(state: AgentState):
    target = state.get("target_language")
    if not target or target == "None":
        return {"status": "No translation needed", "final_report": state.get("final_report", "")}
    
    report = state.get("final_report", "")
    llm = LLMFactory.get_llm(state['llm_provider'], state['llm_model'])
    translator = Agent(
        role='Translator',
        goal=f'Translate the report to {target}',
        backstory="Expert translator.",
        llm=llm
    )
    task = Task(
        description=f"Translate the following report into {target} while preserving formatting:\n\n{report}",
        agent=translator,
        expected_output=f"Report in {target}."
    )
    result = Crew(agents=[translator], tasks=[task]).kickoff()
    return {"final_report": result.raw, "status": f"Translated to {target}"}

def route_quality(state: AgentState):
    score_val = state.get("quality_score", 0)
    score = float(score_val) if score_val is not None else 0.0
    attempts_val = state.get("refinement_attempts", 0)
    attempts = int(attempts_val) if attempts_val is not None else 0
    if score < 1.8 and attempts < 3:
        return "writer"
    return END

# --- GRAPH ---
workflow = StateGraph(AgentState)
workflow.add_node("planner", planning_node)
workflow.add_node("researcher", research_node)
workflow.add_node("writer", writing_node)
workflow.add_node("quality_gate", quality_node)
workflow.add_node("fact_check", fact_check_node)
workflow.add_node("citation", citation_node)
workflow.add_node("translation", translation_node)

workflow.set_entry_point("planner")
workflow.add_edge("planner", "researcher")
workflow.add_edge("researcher", "writer")
workflow.add_edge("writer", "quality_gate")
workflow.add_conditional_edges("quality_gate", route_quality, {"writer": "writer", END: "fact_check"})
workflow.add_edge("fact_check", "citation")
workflow.add_edge("citation", "translation")
workflow.add_edge("translation", END)

app_graph = workflow.compile()

# --- OUTPUT FORMATTING FUNCTIONS (unchanged) ---
def generate_pdf(content: str, topic: str) -> str:
    report_id = hashlib.md5(topic.lower().encode()).hexdigest()[:8]
    filename = f"DeepDive_{report_id}.pdf"
    try:
        doc = SimpleDocTemplate(filename, pagesize=LETTER)
        styles = getSampleStyleSheet()
        story = [
            Paragraph(f"ENTERPRISE RESEARCH: {topic.upper()}", styles['Title']),
            Spacer(1, 24)
        ]
        paragraphs = str(content).split('\n')
        for p in paragraphs:
            p = p.strip()
            if not p:
                continue
            if p.startswith('###'):
                story.append(Paragraph(xml_escape(p.replace('###', '').strip()), styles['Heading3']))
            elif p.startswith('##'):
                story.append(Paragraph(xml_escape(p.replace('##', '').strip()), styles['Heading2']))
            elif p.startswith('#'):
                story.append(Paragraph(xml_escape(p.replace('#', '').strip()), styles['Heading1']))
            elif p.startswith('- ') or p.startswith('* '):
                story.append(Paragraph(f"• {xml_escape(p[2:].strip())}", styles['BodyText']))
            else:
                story.append(Paragraph(xml_escape(p), styles['BodyText']))
            story.append(Spacer(1, 10))
        doc.build(story)
        return filename
    except Exception as e:
        logger.error(f"PDF error: {e}")
        return None

def generate_docx(content: str, topic: str) -> str:
    filename = f"DeepDive_{hashlib.md5(topic.encode()).hexdigest()[:8]}.docx"
    doc = Document()
    doc.add_heading(f"ENTERPRISE RESEARCH: {topic}", 0)
    for para in content.split('\n'):
        if para.strip():
            if para.startswith('#'):
                level = para.count('#')
                doc.add_heading(para.lstrip('#').strip(), level)
            else:
                doc.add_paragraph(para)
    doc.save(filename)
    return filename

def generate_pptx(content: str, topic: str) -> str:
    filename = f"DeepDive_{hashlib.md5(topic.encode()).hexdigest()[:8]}.pptx"
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = f"RESEARCH: {topic}"
    bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
    bullet_slide.shapes.title.text = "Key Takeaways"
    content_shape = bullet_slide.placeholders[1]
    lines = content.split('\n')[:10]
    tf = content_shape.text_frame
    for line in lines:
        if line.strip():
            p = tf.add_paragraph()
            p.text = line.strip()
            p.level = 0
    prs.save(filename)
    return filename

def generate_markdown(content: str, topic: str) -> str:
    filename = f"DeepDive_{hashlib.md5(topic.encode()).hexdigest()[:8]}.md"
    with open(filename, 'w') as f:
        f.write(f"# ENTERPRISE RESEARCH: {topic}\n\n")
        f.write(content)
    return filename

# --- DATABASE SAVE ---
def save_report_to_db(topic: str, email: str, report_text: str, pdf_path: str, user_id: str):
    session = SessionLocal()
    report = ReportDB(topic=topic, email=email, report_text=report_text, pdf_path=pdf_path, user_id=user_id)
    session.add(report)
    session.commit()
    session.close()

# --- EMAIL FUNCTION ---
def send_email(email: str, pdf_path: str, topic: str):
    sender = os.getenv("EMAIL_USER")
    password = os.getenv("EMAIL_PASS")
    if not all([sender, password, pdf_path]):
        logger.error("Missing email config or PDF path")
        return False
    
    msg = MIMEMultipart()
    msg['From'] = sender
    msg['To'] = email
    msg['Subject'] = f"Deep-Dive Research Report: {topic}"
    msg.attach(MIMEText("Please find your comprehensive research report attached.", 'plain'))
    
    try:
        with open(pdf_path, "rb") as f:
            part = MIMEApplication(f.read(), Name=os.path.basename(pdf_path))
            part['Content-Disposition'] = f'attachment; filename="{os.path.basename(pdf_path)}"'
            msg.attach(part)
        
        smtp_server = os.getenv("EMAIL_SMTP_SERVER", "smtp.gmail.com")
        smtp_port = int(os.getenv("EMAIL_SMTP_PORT", 587))
        with smtplib.SMTP(smtp_server, smtp_port) as server:
            server.starttls()
            server.login(sender, password)
            server.send_message(msg)
        logger.info(f"Email sent successfully to {email}")
        return True
    except Exception as e:
        logger.error(f"Email error: {e}")
        return False

# --- HELPER: Extract and render Mermaid diagrams ---
# --- HELPER: Extract and render Mermaid diagrams ---
def render_mermaid_diagrams(text):
    import re
    pattern = r"```mermaid\n(.*?)\n```"
    blocks = re.findall(pattern, text, re.DOTALL)
    
    if not blocks:
        return text
    
    # Remove the code blocks from displayed text
    cleaned = re.sub(pattern, "", text)
    
    # Render each diagram
    for i, block in enumerate(blocks):
        with st.expander(f"Mermaid Diagram {i+1}", expanded=True):
            # Use a cleaner HTML approach with proper initialization
            html_code = f"""
            <div style="background-color: white; padding: 20px; border-radius: 5px;">
                <pre class="mermaid" style="text-align: center; background-color: white;">
                    {block}
                </pre>
            </div>
            <script>
                if (typeof mermaid === 'undefined') {{
                    const script = document.createElement('script');
                    script.src = "https://cdn.jsdelivr.net/npm/mermaid/dist/mermaid.min.js";
                    script.onload = () => {{
                        mermaid.initialize({{
                            startOnLoad: true,
                            theme: 'default',
                            securityLevel: 'loose'
                        }});
                        mermaid.init(undefined, document.querySelectorAll('.mermaid'));
                    }};
                    document.head.appendChild(script);
                }} else {{
                    mermaid.init(undefined, document.querySelectorAll('.mermaid'));
                }}
            </script>
            """
            st.components.v1.html(html_code, height=400)
    
    return cleaned

# --- HELPER: Answer question with report context ---
def answer_question_with_report(question, report, provider, model):
    llm = LLMFactory.get_llm(provider, model, temperature=0.3, max_tokens=1000)
    agent = Agent(
        role='Report Assistant',
        goal='Answer user questions based on the report',
        backstory="You have access to the full report.",
        llm=llm,
        allow_delegation=False
    )
    task = Task(
        description=f"Based on the following report, answer the user's question concisely and accurately.\n\nREPORT:\n{report}\n\nQUESTION: {question}",
        agent=agent,
        expected_output="A helpful answer"
    )
    result = Crew(agents=[agent], tasks=[task]).kickoff()
    return str(result)

# --- HELPER: Generate audio summary ---
def generate_audio_summary(report, provider, model):
    summary = ""
    lines = report.split('\n')
    capturing = False
    for line in lines:
        if "EXECUTIVE SUMMARY" in line.upper():
            capturing = True
            continue
        if capturing and line.strip() == "":
            break
        if capturing:
            summary += line + "\n"
    if not summary:
        summary = ' '.join(report.split()[:200])
    
    audio_bytes = None
    if provider == "OpenAI" and os.getenv("OPENAI_API_KEY"):
        try:
            import openai
            client = openai.OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
            response = client.audio.speech.create(
                model="tts-1",
                voice="nova",
                input=str(summary)[:4096]
            )
            return response.content
        except Exception as e:
            logger.warning(f"OpenAI TTS failed: {e}")
    
    try:
        tts = gTTS(text=str(summary)[:2000], lang='en', slow=False)
        audio_bytes_io = io.BytesIO()
        tts.write_to_fp(audio_bytes_io)
        return audio_bytes_io.getvalue()
    except Exception as e:
        logger.warning(f"gTTS failed: {e}")
    
    return None

# --- SETTINGS PERSISTENCE ---
def update_env_file(key, value):
    env_path = ".env"
    lines = []
    if os.path.exists(env_path):
        with open(env_path, "r") as f:
            lines = f.readlines()
    
    updated = False
    for i, line in enumerate(lines):
        if line.startswith(f"{key}="):
            lines[i] = f"{key}=\"{value}\"\n"
            updated = True
            break
    if not updated:
        lines.append(f"{key}=\"{value}\"\n")
    
    with open(env_path, "w") as f:
        f.writelines(lines)
    os.environ[key] = value

# --- STREAMLIT UI ---
st.set_page_config(page_title="Agentic Studio - Tier 3", layout="wide")

if 'session_id' not in st.session_state:
    st.session_state['session_id'] = hashlib.md5(str(time.time()).encode()).hexdigest()[:12]

# Shared RAG tool per session
if rag_tool is None:
    rag_tool = RAGTool(collection_name=f"docs_{st.session_state['session_id']}")

# Initialize session state for workflow
if 'workflow_step' not in st.session_state:
    st.session_state['workflow_step'] = 'idle'
if 'inputs' not in st.session_state:
    st.session_state['inputs'] = {}
if 'final_state' not in st.session_state:
    st.session_state['final_state'] = None
if 'generated_files' not in st.session_state:
    st.session_state['generated_files'] = {}

# Sidebar menu
with st.sidebar:
    selected = option_menu(
        menu_title="Main Menu",
        options=["New Research", "My Reports", "Settings"],
        icons=["journal-plus", "folder2-open", "gear"],
        menu_icon="cast",
        default_index=0,
    )

if selected == "New Research":
    st.title("AI Agentic Studio")

    # --- FILE UPLOAD (Bring Your Own Data) ---
    with st.sidebar.expander("📁 Upload your own documents", expanded=False):
        uploaded_files = st.file_uploader(
            "Upload PDF, TXT, or CSV files",
            accept_multiple_files=True,
            type=["pdf", "txt", "csv"]
        )
        if uploaded_files and st.button("Process uploaded files"):
            with st.spinner("Adding documents to RAG..."):
                for uploaded_file in uploaded_files:
                    try:
                        text = ""
                        if uploaded_file.type == "application/pdf":
                            pdf_reader = PyPDF2.PdfReader(uploaded_file)
                            for page in pdf_reader.pages:
                                text += page.extract_text()
                        elif uploaded_file.type == "text/plain":
                            text = uploaded_file.read().decode("utf-8")
                        elif uploaded_file.type == "text/csv":
                            csv_content = uploaded_file.read().decode("utf-8")
                            reader = csv.reader(io.StringIO(csv_content))
                            rows = [",".join(row) for row in reader]
                            text = "\n".join(rows)
                        else:
                            st.warning(f"Unsupported file type: {uploaded_file.name}")
                            continue
                        
                        if text:
                            chunks = [text[i:i+1000] for i in range(0, len(text), 1000)]
                            rag_tool.add_documents(chunks, metadatas=[{"source": uploaded_file.name}] * len(chunks))
                            st.success(f"Added {uploaded_file.name}")
                    except Exception as e:
                        st.error(f"Error processing {uploaded_file.name}: {e}")

    with st.form("research_form"):
        col1, col2 = st.columns(2)
        with col1:
            topic = st.text_input("Research Topic", "Impact of GenAI on Software Development Lifecycle 2026")
            email = st.text_input("Email (optional)")
            num_agents = st.slider("Number of Researchers", 1, 5, 3)
            writing_style = st.selectbox("Writing Style", ["formal", "persuasive", "technical"])
        with col2:
            llm_provider = st.selectbox("LLM Provider", ["OpenAI", "Anthropic", "Gemini", "Ollama"])
            # Dynamically set model options based on provider
            model_options = FREE_MODELS[llm_provider]
            llm_model = st.selectbox("Model", model_options, index=0)
            target_language = st.selectbox("Translate to (optional)", ["None", "Spanish", "French", "German", "Chinese"])
            edit_plan = st.checkbox("Edit research plan before execution")

        submitted = st.form_submit_button("🚀 Launch Deep Research")

    if submitted:
        # --- API Key Validation ---
        valid = True
        if llm_provider == "OpenAI" and not os.getenv("OPENAI_API_KEY"):
            st.error("OpenAI API key is missing. Please set it in Settings.")
            valid = False
        if llm_provider == "Anthropic" and not os.getenv("ANTHROPIC_API_KEY"):
            st.error("Anthropic API key is missing. Please set it in Settings.")
            valid = False
        if llm_provider == "Gemini" and not os.getenv("GOOGLE_API_KEY"):
            st.error("Google API key is missing. Please set it in Settings.")
            valid = False
        if llm_provider == "Ollama":
            ollama_url = os.getenv("OLLAMA_URL", "http://localhost:11434")
            try:
                response = requests.get(f"{ollama_url}/api/tags", timeout=5)
                if response.status_code == 200:
                    models = response.json().get('models', [])
                    available_models = [m['name'] for m in models]
                    if not any(llm_model.lower() in m.lower() for m in available_models):
                        st.error(f"❌ Model '{llm_model}' not found in Ollama.")
                        valid = False
                else:
                    st.error(f"Ollama returned error {response.status_code}")
                    valid = False
            except requests.exceptions.ConnectionError:
                st.error(f"❌ Cannot connect to Ollama at {ollama_url}.")
                valid = False
        
        if valid:
            inputs = {
                "topic": topic,
                "email": email,
                "num_agents": num_agents,
                "writing_style": writing_style,
                "llm_provider": llm_provider,
                "llm_model": llm_model,
                "target_language": None if target_language=="None" else target_language,
                "plan": "",
                "raw_research": "",
                "final_report": "",
                "quality_score": 0.0,
                "status": "Starting",
                "citations": [],
                "research_notes": [],
                "fact_check_report": "",
                "readability_scores": {},
                "cache_key": hashlib.md5(topic.encode()).hexdigest(),
                "refinement_attempts": 0
            }
            st.session_state['inputs'] = inputs
            st.session_state['generated_files'] = {} # Reset files
            
            if edit_plan:
                with st.spinner("Generating plan..."):
                    planner_out = planning_node(inputs)
                    st.session_state['inputs']['plan'] = planner_out['plan']
                    st.session_state['workflow_step'] = 'edit_plan'
            else:
                st.session_state['workflow_step'] = 'execute'
            st.rerun()

    # --- Workflow Execution Logic ---
    if st.session_state['workflow_step'] == 'edit_plan':
        st.subheader("📋 Review Research Plan")
        edited_plan = st.text_area("You can modify the research structure before agents begin:", 
                                   value=st.session_state['inputs']['plan'], height=300)
        if st.button("🚀 Start Research with this Plan"):
            st.session_state['inputs']['plan'] = edited_plan
            st.session_state['workflow_step'] = 'execute'
            st.rerun()

    if st.session_state['workflow_step'] == 'execute':
        st.session_state['start_time'] = time.time()
        final_state = {}
        with st.status("🤖 Deep Research Agents at work...", expanded=True) as status:
            try:
                # No interrupt used here for simplicity as we handle planning before execution
                for output in app_graph.stream(st.session_state['inputs']):
                    for node, data in output.items():
                        if data is None:
                            st.warning(f"⚠️ **{node}** returned no data. Skipping.")
                            continue
                        status_text = data.get('status', 'Done') if isinstance(data, dict) else 'Done'
                        st.write(f"✅ **{node}**: {status_text}")
                        if isinstance(data, dict):
                            final_state.update(data)
                
                duration = round(time.time() - st.session_state['start_time'], 2)
                st.session_state['duration'] = duration
                status.update(label=f"✅ Deep Dive Complete in {duration}s!", state="complete")
                st.session_state['final_state'] = final_state
                st.session_state['workflow_step'] = 'results'
                st.rerun()
            except Exception as e:
                st.exception(e)
                status.update(label="❌ Failed", state="error")
                st.session_state['workflow_step'] = 'idle'
                st.stop()

    if st.session_state['workflow_step'] == 'results' and st.session_state['final_state']:
        final_state = st.session_state['final_state']
        report_text = final_state.get("final_report", "")

        # Render Mermaid diagrams
        cleaned_report = render_mermaid_diagrams(report_text)

        col1, col2 = st.columns([2, 1])

        with col1:
            st.subheader("📄 Full Executive Whitepaper")
            st.markdown(cleaned_report)

            # Chat with your report
            st.subheader("💬 Chat with your report")
            if "messages" not in st.session_state:
                st.session_state.messages = []
            for msg in st.session_state.messages:
                with st.chat_message(msg["role"]):
                    st.markdown(msg["content"])
            if prompt := st.chat_input("Ask about the report..."):
                st.session_state.messages.append({"role": "user", "content": prompt})
                with st.chat_message("user"):
                    st.markdown(prompt)
                with st.chat_message("assistant"):
                    with st.spinner("Thinking..."):
                        answer = answer_question_with_report(
                            prompt, 
                            report_text, 
                            final_state.get('llm_provider', 'OpenAI'),
                            final_state.get('llm_model', 'gpt-3.5-turbo')
                        )
                    st.markdown(answer)
                st.session_state.messages.append({"role": "assistant", "content": answer})

            # Feedback loop
            feedback = st.radio("Was this report helpful?", ("Yes", "No"), key="feedback")
            if feedback == "No":
                correction = st.text_area("Please suggest corrections or additional points:")
                if st.button("Submit Feedback & Refine"):
                    st.info("Refinement triggered (simulated).")

        with col2:
            st.subheader("📦 Professional Deliverables")
            
            # Use cached file paths
            if 'pdf' not in st.session_state['generated_files']:
                st.session_state['generated_files']['pdf'] = generate_pdf(report_text, topic)
            if 'docx' not in st.session_state['generated_files']:
                st.session_state['generated_files']['docx'] = generate_docx(report_text, topic)
            if 'pptx' not in st.session_state['generated_files']:
                st.session_state['generated_files']['pptx'] = generate_pptx(report_text, topic)
            if 'md' not in st.session_state['generated_files']:
                st.session_state['generated_files']['md'] = generate_markdown(report_text, topic)

            pdf_path = st.session_state['generated_files']['pdf']
            docx_path = st.session_state['generated_files']['docx']
            pptx_path = st.session_state['generated_files']['pptx']
            md_path = st.session_state['generated_files']['md']

            if pdf_path:
                with open(pdf_path, "rb") as f:
                    st.download_button("💾 Download PDF", f, file_name=pdf_path)
            if docx_path:
                with open(docx_path, "rb") as f:
                    st.download_button("📄 Download DOCX", f, file_name=docx_path)
            if pptx_path:
                with open(pptx_path, "rb") as f:
                    st.download_button("📽️ Download PPTX", f, file_name=pptx_path)
            if md_path:
                with open(md_path, "rb") as f:
                    st.download_button("📝 Download Markdown", f, file_name=md_path)

            if email and st.button("📧 Send to Email"):
                with st.spinner("Dispatching..."):
                    if send_email(email, pdf_path, topic):
                        st.success("✅ Dispatched via SMTP!")
                    else:
                        st.warning("Check .env credentials")

            if st.button("💾 Save to My Reports"):
                save_report_to_db(topic, email, report_text, pdf_path, "guest_user")
                st.success("Report saved to database!")

            # Audio Brief
            st.subheader("🎧 Audio Brief")
            if st.button("Generate Audio Summary"):
                with st.spinner("Creating audio..."):
                    audio_bytes = generate_audio_summary(report_text, final_state.get('llm_provider'), final_state.get('llm_model'))
                if audio_bytes:
                    st.audio(audio_bytes, format="audio/mp3")
                else:
                    st.error("Could not generate audio.")

            # Source Tracker
            st.subheader("🔗 Sources")
            citations = final_state.get("citations", [])
            if citations:
                with st.expander("View citations", expanded=False):
                    for i, cit in enumerate(citations):
                        st.markdown(f"{i+1}. [{cit['source']}]({cit['source']}) - {cit['snippet'][:100]}...")
            else:
                st.info("No citations available.")

            readability = final_state.get("readability_scores", {})
            st.json({
                "execution_time": f"{st.session_state.get('duration', 0)}s",
                "quality_depth_score": round(float(final_state.get("quality_score", 0)), 2) if final_state.get("quality_score") is not None else 0.0,
                "flesch_reading_ease": readability.get("flesch_reading_ease", "N/A"),
                "gunning_fog": readability.get("gunning_fog", "N/A"),
                "complex_sentences": readability.get("complex_sentences_count", 0),
                "fact_check": final_state.get("fact_check_report", "Not performed")[:200] + "..."
            })

elif selected == "My Reports":
    st.title("📚 My Saved Reports")
    session = SessionLocal()
    reports = session.query(ReportDB).all()
    for r in reports:
        with st.expander(f"{r.topic} - {r.created_at.strftime('%Y-%m-%d %H:%M')}"):
            st.markdown(r.report_text[:500] + "...")
            if r.pdf_path and os.path.exists(r.pdf_path):
                with open(r.pdf_path, "rb") as f:
                    st.download_button("Download PDF", f, file_name=r.pdf_path)
    session.close()

elif selected == "Settings":
    st.title("⚙️ Settings")
    st.write("Configure your preferences and integrations.")
    with st.expander("API Keys"):
        openai_key = st.text_input("OpenAI API Key", type="password", value=os.getenv("OPENAI_API_KEY", ""))
        if st.button("Save OpenAI Key"):
            update_env_file("OPENAI_API_KEY", openai_key)
            st.success("Saved!")
            
        anthropic_key = st.text_input("Anthropic API Key", type="password", value=os.getenv("ANTHROPIC_API_KEY", ""))
        if st.button("Save Anthropic Key"):
            update_env_file("ANTHROPIC_API_KEY", anthropic_key)
            st.success("Saved!")
            
        google_key = st.text_input("Google API Key", type="password", value=os.getenv("GOOGLE_API_KEY", ""))
        if st.button("Save Google Key"):
            update_env_file("GOOGLE_API_KEY", google_key)
            st.success("Saved!")
            
        news_key = st.text_input("News API Key", type="password", value=os.getenv("NEWS_API_KEY", ""))
        if st.button("Save News API Key"):
            update_env_file("NEWS_API_KEY", news_key)
            st.success("Saved!")
            
    with st.expander("Email"):
        email_user = st.text_input("Email User", value=os.getenv("EMAIL_USER", ""))
        email_pass = st.text_input("Email Password", type="password", value=os.getenv("EMAIL_PASS", ""))
        smtp_server = st.text_input("SMTP Server", value=os.getenv("EMAIL_SMTP_SERVER", "smtp.gmail.com"))
        smtp_port = st.text_input("SMTP Port", value=os.getenv("EMAIL_SMTP_PORT", "587"))
        
        if st.button("Save Email Settings"):
            update_env_file("EMAIL_USER", email_user)
            update_env_file("EMAIL_PASS", email_pass)
            update_env_file("EMAIL_SMTP_SERVER", smtp_server)
            update_env_file("EMAIL_SMTP_PORT", smtp_port)
            st.success("Saved!")
    with st.expander("Monitoring"):
        enable_monitoring = st.checkbox("Enable Prometheus metrics", value=os.getenv("ENABLE_MONITORING", "false")=="true")
        if enable_monitoring:
            os.environ["ENABLE_MONITORING"] = "true"
        else:
            os.environ["ENABLE_MONITORING"] = "false"
    st.info("Changes take effect immediately for environment variables; some may require app restart.")